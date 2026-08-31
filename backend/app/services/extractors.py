import re
from pathlib import Path

IMAGE_EXTENSIONS = {'.jpg', '.jpeg', '.png', '.webp', '.gif', '.bmp', '.tiff'}
DOCUMENT_EXTENSIONS = {'.txt', '.pdf', '.docx', '.pptx'}
def category_for(path: Path) -> str: return 'image' if path.suffix.lower() in IMAGE_EXTENSIONS else 'document' if path.suffix.lower() in DOCUMENT_EXTENSIONS else 'other'

def normalize_text(text: str) -> str:
    text = text.replace('\u00ad', '').replace('\r\n', '\n')
    text = re.sub(r'[ \t]+', ' ', text)
    return re.sub(r'\n{3,}', '\n\n', text).strip().lower()

def extract_text(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == '.txt': return normalize_text(path.read_text(encoding='utf-8', errors='ignore'))
    if suffix == '.pdf':
        import fitz
        with fitz.open(path) as doc: return normalize_text('\n'.join(page.get_text() for page in doc))
    if suffix == '.docx':
        from docx import Document
        return normalize_text('\n'.join(p.text for p in Document(path).paragraphs))
    if suffix == '.pptx':
        from pptx import Presentation
        prs = Presentation(path); return normalize_text('\n'.join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, 'text')))
    return ''

def image_metadata(path: Path) -> dict:
    from PIL import Image
    import imagehash
    with Image.open(path) as image:
        return {'dimensions': f'{image.width} × {image.height}', 'perceptual_hashes': {'phash': str(imagehash.phash(image)), 'dhash': str(imagehash.dhash(image)), 'ahash': str(imagehash.average_hash(image))}}
